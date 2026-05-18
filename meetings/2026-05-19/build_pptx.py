#!/usr/bin/env python3
"""
Build polished SRIP_Presentation_v2.pptx
- Preserves the original theme/layout from the uploaded file
- Rewrites all 9 slides + adds 5 new ones = 14 slides total
- Matches exact XML formatting from original (17.3pt, no bullets, 95% line-spacing)
- Embeds demo images on relevant slides
"""

import zipfile, copy, re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor

INPUT   = '/root/.claude/uploads/6bae68ba-639c-445a-9fcf-954d530a3ca4/f4789e80-SRIP_Presena.pptx'
OUTPUT  = '/home/user/ECESRIP/meetings/2026-05-19/SRIP_Presentation_v2.pptx'
DEMO    = Path('/home/user/ECESRIP/meetings/2026-05-11/demo')

# ── EMU constants (from original slide geometry) ──────────────────────────────
TITLE_L, TITLE_T, TITLE_W, TITLE_H = 311700, 445025,  8520600, 572700
BODY_L,  BODY_T,  BODY_W,  BODY_H  = 311700, 1152475, 8520600, 3416400

# For split-layout slides (text left, image right)
SPLIT_TEXT_L, SPLIT_TEXT_T = 311700, 1152475
SPLIT_TEXT_W, SPLIT_TEXT_H = 3400000, 3500000

# ── XML namespaces ─────────────────────────────────────────────────────────────
A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def _x(tag, ns=A):
    return f'{{{ns}}}{tag}'


# ── XML builder helpers ────────────────────────────────────────────────────────

def para_xml(text, sz=1730, bold=False, color='000000', space_before=0,
             line_pct=95, indent=None, italic=False):
    """One <a:p> element matching the original slide style."""
    b_attr   = ' b="1"' if bold else ''
    i_attr   = ' i="1"' if italic else ''
    ind_attr = f' indent="{indent}"' if indent is not None else ' indent="0"'
    escaped  = (text.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;'))
    xml = (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr{ind_attr} lvl="0" marL="0" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="{line_pct}000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="{space_before}"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buNone/></a:pPr>'
        f'<a:r><a:rPr lang="en" sz="{sz}"{b_attr}{i_attr}>'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'</a:rPr><a:t>{escaped}</a:t></a:r>'
        f'<a:endParaRPr sz="{sz}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:endParaRPr>'
        f'</a:p>'
    )
    return etree.fromstring(xml)

def blank_para():
    xml = (f'<a:p xmlns:a="{A}"><a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
           f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
           f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
           f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
           f'<a:buNone/></a:pPr>'
           f'<a:r><a:t></a:t></a:r>'
           f'<a:endParaRPr sz="1730"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:endParaRPr>'
           f'</a:p>')
    return etree.fromstring(xml)

def shape_xml(shape_id, ph_type, ph_idx, left, top, w, h,
              title_text=None, body_paras=None,
              title_color='158158', body_font_sz=1730,
              body_anchor='t', body_auto='normAutofit'):
    """Full <p:sp> element for either a title or body placeholder."""
    if ph_type == 'title':
        ph_xml  = f'<p:ph type="title"/>'
        name    = f'Google Shape;{shape_id};title'
        anchor  = 'b'
        auto    = f'<a:normAutofit fontScale="90000"/>'
        body_pr = (f'<a:bodyPr anchorCtr="0" anchor="{anchor}" bIns="91425" lIns="91425" '
                   f'spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">{auto}</a:bodyPr>')
        if title_text:
            t_esc = title_text.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
            paras_xml = (
                f'<a:p xmlns:a="{A}"><a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
                f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
                f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
                f'<a:buNone/></a:pPr>'
                f'<a:r><a:rPr lang="en">'
                f'<a:solidFill><a:srgbClr val="{title_color}"/></a:solidFill>'
                f'</a:rPr><a:t>{t_esc}</a:t></a:r><a:endParaRPr/></a:p>'
            )
        else:
            paras_xml = ''
    else:  # body
        ph_xml  = f'<p:ph idx="1" type="body"/>'
        name    = f'Google Shape;{shape_id};body'
        auto    = '' if body_auto == 'normAutofit' else ''
        body_pr = (f'<a:bodyPr anchorCtr="0" anchor="{body_anchor}" bIns="91425" lIns="91425" '
                   f'spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">'
                   f'<a:normAutofit/></a:bodyPr>')
        if body_paras:
            paras_xml = '\n'.join(etree.tostring(p, encoding='unicode') for p in body_paras)
        else:
            paras_xml = ''

    xml = (
        f'<p:sp xmlns:p="{P}" xmlns:a="{A}">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{shape_id}" name="{name}"/>'
        f'<p:cNvSpPr txBox="1"/>'
        f'<p:nvPr>{ph_xml}</p:nvPr>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'</p:spPr>'
        f'<p:txBody>{body_pr}<a:lstStyle/>{paras_xml}</p:txBody>'
        f'</p:sp>'
    )
    return etree.fromstring(xml)


# ── Slide content definitions ──────────────────────────────────────────────────
# Each entry: (title, [(text, bold, color, space_before), ...], image_path_or_None, split_layout)
# image_path: image placed on right half if split_layout=True, else full-width below title
# color codes: '' means inherit (title color); '000000' = black; '158158' = green; '058DC7' = blue

GREEN = '158158'
BLACK = '000000'
BLUE  = '058DC7'
DGREY = '444444'

def B(t, sb=0):  return (t, True,  BLACK, sb)
def N(t, sb=0):  return (t, False, BLACK, sb)
def G(t, sb=0):  return (t, True,  GREEN, sb)
def Bl(t, sb=0): return (t, False, BLUE,  sb)
def empty(sb=0): return ('', False, BLACK, sb)

SLIDES = [
    # ── 1. Title (keep original — signal with None body) ─────────────────────
    None,

    # ── 2. Hopfield Neural Networks ──────────────────────────────────────────
    ("Hopfield Neural Networks", [
        N("Fully connected recurrent network — N bipolar neurons, sᵢ ∈ {-1, +1}"),
        N("Symmetric weight matrix W with zero diagonal — no self-connections"),
        N("Training writes W once; stores M patterns as energy minima"),
        N("Energy:  E(s) = −½ sᵀWs  — stored patterns sit at local minima"),
        N("Inference: start from corrupted input, converge to nearest stored memory"),
        N("Unlike deep nets — no layers, no backprop, no gradient descent at inference"),
    ], None, False),

    # ── 3. Hebbian Learning ───────────────────────────────────────────────────
    ("Hebbian Learning", [
        G("Wᵢⱼ = (1/N) Σ_μ  ξᵢ^μ ξⱼ^μ ,    Wᵢᵢ = 0", sb=200),
        empty(),
        N("Neurons that agree in a pattern → positive weight; disagree → negative"),
        N("All M patterns accumulated in one shot — no iterations, order-independent"),
        N("Local: Wᵢⱼ depends only on neurons i and j across patterns"),
        N("Incremental: new pattern = rank-1 update ΔW = ξξᵀ / N"),
        N("Storage capacity ≈ 0.138N — spurious states dominate above this"),
        N("Note: inverted patterns (−ξ) are also attractors"),
    ], None, False),

    # ── 4. Storkey Learning ───────────────────────────────────────────────────
    ("Storkey Learning", [
        G("ΔW = (1/N)(ξ^μ(ξ^μ)ᵀ − h(ξ^μ)ᵀ − ξ^μ hᵀ),   hᵢ = Σⱼ≠ᵢ Wᵢⱼ ξⱼ^μ", sb=200),
        empty(),
        N("hᵢ = local field — what current weights already predict for pattern μ"),
        N("Correction terms subtract existing crosstalk before writing the new pattern"),
        N("Updated once per pattern, in presentation order — earlier patterns may be overwritten"),
        N("Local and incremental: bounded work per pattern, no global recomputation"),
        N("Storage capacity ≈ 0.14–0.22N — significantly better than Hebbian"),
        N("~2× compute per pattern vs Hebbian, but far fewer spurious states above capacity"),
    ], None, False),

    # ── 5. Pseudo-inverse Learning ───────────────────────────────────────────
    ("Pseudo-inverse Learning", [
        G("W = Ξ (ΞᵀΞ)⁻¹ Ξᵀ    where Ξ = [ξ¹|…|ξᴹ]  (N×M pattern matrix)", sb=200),
        empty(),
        N("Theoretically optimal — stores up to N linearly independent patterns"),
        N("Weights W form the orthogonal projection onto the pattern subspace"),
        B("NOT incremental: one new pattern → full pseudo-inverse recomputed over all M patterns"),
        B("NOT local: each Wᵢⱼ depends on all M patterns simultaneously via matrix inversion"),
        N("Hardware cost: O(N²M + M³) matrix inversion — no online update circuit is feasible"),
        N("For our project (offline training → static LUT): technically usable, but Storkey preferred"),
    ], None, False),

    # ── 6. Local & Incremental ───────────────────────────────────────────────
    ("Local & Incremental — Why It Matters", [
        B("Local:", sb=100),
        N("Wᵢⱼ is determined only by neurons i and j — no global communication during training"),
        N("Maps directly to hardware: each weight update is an independent multiply-accumulate"),
        empty(),
        B("Incremental:", sb=100),
        N("Add patterns one at a time — bounded, fixed work per pattern (rank-1 update)"),
        N("Hardware relearning: Storkey → small fixed update circuit"),
        N("Hardware relearning: pseudo-inverse → full matrix inverter on chip — impractical"),
        empty(100),
        N("Pseudo-inverse: global (all patterns coupled) + batch (all patterns needed at once)"),
    ], None, False),

    # ── 7. Inference Process ─────────────────────────────────────────────────
    ("Inference Process", [
        G("sᵢ ← sign( hᵢ ),    hᵢ = Σⱼ Wᵢⱼ sⱼ   (local field — weighted vote)", sb=200),
        empty(),
        N("Each neuron checks its neighbours: positive field → +1, negative field → −1"),
        N("Tie (hᵢ = 0): hold current state — no change"),
        N("Sweep through all neurons repeatedly until no neuron changes state"),
        N("Fixed point = stored memory (perfect recall) or spurious state (above capacity)"),
        empty(100),
        N("LUT insight: sign(hᵢ) is a Boolean function of N binary inputs → truth table"),
        N("Only the sign matters, not the magnitude — hardware needs no multiplier"),
    ], None, False),

    # ── 8. Update Order Semantics ─────────────────────────────────────────────
    ("Update Order Semantics", [
        B("Async-cyclic (chosen for hardware):"),
        N("One neuron per clock cycle, fixed order 0→1→…→N−1, repeat. Fully deterministic"),
        empty(),
        B("Async-random:"),
        N("Random order each sweep — unbiased, same convergence guarantee as cyclic"),
        empty(),
        B("Sync:"),
        N("All N neurons update simultaneously — fast throughput but unsafe (2-cycles possible)"),
        empty(100),
        N("Convergence proof (async only): each flip ΔE = −hᵢΔsᵢ ≤ 0 — energy never increases"),
        N("Goles & Olivos (1980): sync with symmetric W → fixed point OR infinite 2-cycle oscillation"),
    ], None, False),

    # ── 9. The 2-Cycle Problem ───────────────────────────────────────────────
    ("The 2-Cycle Problem (Sync Update)", [
        N("2-cycle: states A and B where  sync(A) = B  and  sync(B) = A  — infinite oscillation"),
        N("Same weights, same start — only update mode differs (top row sync, bottom async-cyclic)"),
        N("Detection rule:  s[t] == s[t−2]  AND  s[t] ≠ s[t−1]"),
        N("Not an edge case — demonstrated with M=4 patterns in N=16 above capacity (load = 0.25)"),
        N("Hardware sync implementation must include a 2-cycle detector to guarantee termination"),
    ], 'demo_sync_2cycle.png', True),

    # ── 10. Training Code Capabilities ───────────────────────────────────────
    ("Training Code Capabilities", [
        B("Learning rules — one-line switch:"),
        N("RULE = STORKEY  /  HEBBIAN   (sim/python/hopfield_net.py, line 20)"),
        empty(),
        B("Update modes — one-line switch:"),
        N("UPDATE_MODE = ASYNC_CYCLIC  /  ASYNC_RANDOM  /  SYNC"),
        empty(),
        B("Benchmark sweep  (sim/python/benchmark.py):"),
        N("N ∈ {6,8,10,12,14,16}, load M/N up to 0.30, 9 noise levels, 50 trials per point"),
        N("Metrics: recall accuracy, basin width, spurious rate, mean convergence sweeps"),
        N("Output: timestamped CSV → sim/results/,  plots → sim/plots/"),
    ], None, False),

    # ── 11. Datasets & Benchmarks ────────────────────────────────────────────
    ("Datasets & Benchmarks", [
        B("Random binary (built-in):"),
        N("Synthetic N-bit patterns, M configurable — fully reproducible, no external data"),
        empty(),
        B("sklearn 8×8 digits (built-in, no internet):"),
        N("1,797 hand-written digit images (0–9), binarised at threshold 8 → {−1, +1}"),
        empty(),
        B("Key results from benchmark sweep:"),
        N("Both rules: zero spurious rate below 0.138N — matches theoretical capacity exactly"),
        N("Above capacity (N=8, load=0.38): Hebbian → 79% spurious;  Storkey → ~0%"),
        N("Basin width scales with N:  N=16, M=1 → 6 bits;   N=6, M=1 → 2 bits"),
    ], None, False),

    # ── 12. LUT Synthesis Pipeline ────────────────────────────────────────────
    ("LUT Synthesis Pipeline", [
        G("W (N×N weights, Python)"),
        Bl("    ↓  truth_table_gen.py  →  2ᴺ rows per neuron  (.pla format)"),
        Bl("    ↓  Espresso minimiser   →  minimal Sum-of-Products Boolean expression"),
        Bl("    ↓  sv_emitter.py        →  SystemVerilog  always_comb  blocks"),
        Bl("    ↓  Quartus synthesis    →  LUT-mapped netlist on Cyclone V FPGA"),
        empty(200),
        N("Hazard-free mode: Espresso -Dhazard adds consensus terms — safe for single-variable transitions"),
        N("Feasibility:  N≤10 comfortable (~1K rows/neuron),  N=16 borderline (~64K rows)"),
        N("N>16 → sparse: keep F strongest weights per neuron → 2ᶠ rows  (LogicNets approach)"),
    ], None, False),

    # ── 13. Research Roadmap ─────────────────────────────────────────────────
    ("Research Roadmap", [
        G("Phase 1  (next step): Clocked Baseline", sb=100),
        N("Python training → truth tables → Espresso → clocked SystemVerilog → ModelSim + Quartus"),
        N("One flip-flop per neuron, one LUT evaluation per clock cycle, counter FSM"),
        N("Goal: formal termination proof by energy argument, ModelSim RTL vs Python agreement"),
        empty(150),
        G("Phase 2: Async Combinational Feedback"),
        N("Strip flip-flops — wire LUT outputs directly back to inputs"),
        N("Demonstrate convergence without clock; analyse hazard-free correctness"),
        empty(150),
        G("Phase 3: Ising Machine"),
        N("Encode NP-hard problem instance as weight matrix W — run hardware to find ground state"),
        N("Approximate combinatorial optimisation using the same LUT synthesis pipeline"),
    ], None, False),

    # ── 14. Demo: Hopfield Recall ─────────────────────────────────────────────
    ("Demo: Hopfield Recall", [
        N("N=64, M=8 patterns  (load = 0.125, below 0.138N capacity)"),
        N("sklearn 8×8 digits 0–7, binarised — 25% pixel noise, Storkey, async-cyclic"),
        N("All 8 patterns recalled perfectly in ≤ 2 sweeps = ≤ 128 clock cycles"),
        N("Energy curve per neuron-update step — strictly decreasing to fixed point"),
    ], 'demo_recall_grid.png', True),
]


# ── Build the presentation ────────────────────────────────────────────────────
prs = Presentation(INPUT)

# ── Delete slides 2–9 (keep slide 0 = title) ─────────────────────────────────
from pptx.oxml.ns import qn as _qn

def _delete_slide(prs, idx):
    """Remove slide at index idx from the presentation."""
    slides_elem = prs.slides._sldIdLst
    # find rId for this slide
    slide_part = prs.slides._sldIdLst[idx]
    rId = slide_part.get(_qn('r:id'))
    slides_elem.remove(slide_part)
    # drop the relationship so python-pptx doesn't choke
    prs.slides.part.drop_rel(rId)

# delete from the end to avoid index shifting
for i in range(len(prs.slides) - 1, 0, -1):
    _delete_slide(prs, i)

# ── Helper to add a new slide from the body layout ────────────────────────────
def get_layout(name_fragment):
    for lay in prs.slide_layouts:
        if name_fragment.upper() in lay.name.upper():
            return lay
    return prs.slide_layouts[1]   # fallback

body_layout  = get_layout('BODY')
title_layout = get_layout('TITLE')

_shape_counter = [100]

def next_id():
    _shape_counter[0] += 1
    return _shape_counter[0]

# ── Add slides ────────────────────────────────────────────────────────────────
for slide_def in SLIDES[1:]:   # skip first None (title slide kept)
    title_text, bullets, img_file, split = slide_def

    slide = prs.slides.add_slide(body_layout)

    # Remove all placeholder shapes added by the layout
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        if sp.tag not in (_qn('p:sp'), _qn('p:nvGrpSpPr'), _qn('p:grpSpPr')):
            continue
        if sp.tag == _qn('p:sp'):
            sp_tree.remove(sp)

    # ── Title shape ───────────────────────────────────────────────────────────
    t_shape = shape_xml(next_id(), 'title', None,
                        TITLE_L, TITLE_T, TITLE_W, TITLE_H,
                        title_text=title_text)
    sp_tree.append(t_shape)

    # ── Body layout: split (text left + image right) or full-width text ────────
    if split and img_file and (DEMO / img_file).exists():
        # Narrow text column on the left
        body_paras = [blank_para()]
        for text, bold, color, sb in bullets:
            if text == '':
                body_paras.append(blank_para())
            else:
                body_paras.append(para_xml(text, bold=bold, color=color, space_before=sb,
                                           sz=1500))   # slightly smaller for split layout
        b_shape = shape_xml(next_id(), 'body', 1,
                            SPLIT_TEXT_L, SPLIT_TEXT_T,
                            SPLIT_TEXT_W, SPLIT_TEXT_H,
                            body_paras=body_paras)
        sp_tree.append(b_shape)

        # Image on the right
        from PIL import Image as PILImg
        img_path = DEMO / img_file
        im = PILImg.open(img_path)
        pw, ph = im.size   # pixels

        slide_w = int(prs.slide_width)    # EMU
        slide_h = int(prs.slide_height)   # EMU

        # Right column: from x=3.9in to slide right edge
        img_left = Emu(3_566_160)         # 3.9 inches
        img_top  = Emu(1_097_280)         # 1.2 inches (just below title)
        avail_w  = slide_w - int(img_left) - 182880   # leave 0.2in margin
        avail_h  = slide_h - int(img_top)  - 182880

        aspect = ph / pw
        img_w = avail_w
        img_h = int(img_w * aspect)
        if img_h > avail_h:
            img_h = avail_h
            img_w = int(img_h / aspect)

        slide.shapes.add_picture(str(img_path),
                                 img_left, img_top,
                                 width=Emu(img_w), height=Emu(img_h))
    else:
        # Full-width text body
        body_paras = [blank_para()]
        for text, bold, color, sb in bullets:
            if text == '':
                body_paras.append(blank_para())
            else:
                body_paras.append(para_xml(text, bold=bold, color=color, space_before=sb))
        b_shape = shape_xml(next_id(), 'body', 1,
                            BODY_L, BODY_T, BODY_W, BODY_H,
                            body_paras=body_paras)
        sp_tree.append(b_shape)

prs.save(OUTPUT)
print(f"Saved {len(prs.slides)} slides → {OUTPUT}")
