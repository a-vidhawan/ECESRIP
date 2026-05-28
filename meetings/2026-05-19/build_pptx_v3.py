"""
Build SRIP_Presentation_v3.pptx
Topic: HNN as Ising Machines & Gradient-Based Training Methods
Theme: cloned exactly from 0b712d9c-SRIP_Presena.pptx
"""

import copy, re, shutil
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

INPUT  = '/root/.claude/uploads/b83f060a-0718-412b-ba15-7f9efadcb4bb/0b712d9c-SRIP_Presena.pptx'
OUTPUT = '/home/user/ECESRIP/meetings/2026-05-19/SRIP_Presentation_v3.pptx'

A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── helpers to build paragraph XML ───────────────────────────────────────────

def _rpr(bold=False):
    return (f'<a:rPr xmlns:a="{A}" lang="en" sz="1730"'
            + (' b="1"' if bold else '')
            + f'><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:rPr>')

def _endParaRpr():
    return f'<a:endParaRPr xmlns:a="{A}" sz="1730"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:endParaRPr>'

def plain_para(text, bold=False):
    """Plain (no-bullet) paragraph at sz=1730, 95% line spacing."""
    escaped = text.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buNone/></a:pPr>'
        f'<a:r>{_rpr(bold)}<a:t>{escaped}</a:t></a:r>'
        f'{_endParaRpr()}</a:p>'
    )

def bullet_para(text, bold_part=None):
    """Bullet paragraph (● char) at sz=1730."""
    if bold_part and bold_part in text:
        before, after = text.split(bold_part, 1)
        be = before.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
        bo = bold_part.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
        af = after.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
        runs = (f'<a:r>{_rpr()}<a:t>{be}</a:t></a:r>'
                f'<a:r>{_rpr(bold=True)}<a:t>{bo}</a:t></a:r>'
                f'<a:r>{_rpr()}<a:t>{af}</a:t></a:r>')
    else:
        escaped = text.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;')
        runs = f'<a:r>{_rpr()}<a:t>{escaped}</a:t></a:r>'
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="-338455" lvl="0" marL="457200" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buClr><a:srgbClr val="000000"/></a:buClr>'
        f'<a:buSzPts val="1730"/><a:buChar char="&#x25CF;"/></a:pPr>'
        f'{runs}'
        f'{_endParaRpr()}</a:p>'
    )

def empty_para():
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buNone/></a:pPr>'
        f'<a:r>{_rpr()}<a:t></a:t></a:r>'
        f'{_endParaRpr()}</a:p>'
    )

def P(label):  return ('p', label)   # plain
def B(label):  return ('b', label)   # bullet
def E():       return ('e', '')      # empty spacer

# ── slide content ─────────────────────────────────────────────────────────────

SLIDES = [
    # (title, [(type, text), ...])
    (
        None,   # title slide — handled separately
        []
    ),
    (
        "Ising Machine — Connection to HNNs",
        [
            P("The Ising Hamiltonian H = −Σ Jᵢⱼσᵢσⱼ − Σ hᵢσᵢ maps exactly onto Hopfield energy E = −½ sᵀWs − bᵀs"),
            E(),
            B("Set Wᵢⱼ = 2Jᵢⱼ and bᵢ = hᵢ — the async update sᵢ ← sign(Wᵢ·s + bᵢ) minimises H"),
            B("Finding the ground state (minimum H) is NP-hard — a converging HNN IS an Ising machine solver"),
            B("Symmetry of J required for convergence guarantee — already satisfied by Storkey / Hebbian W"),
            E(),
            P("Our code is already a valid Ising machine — one missing piece: bias vector b"),
        ]
    ),
    (
        "QUBO ↔ Ising Conversion",
        [
            P("Binary {0,1} variables xᵢ relate to bipolar spins via xᵢ = (σᵢ + 1)/2"),
            P("Any QUBO  min xᵀQx  rewrites in Ising form as:"),
            E(),
            B("Jᵢⱼ = −Qᵢⱼ / 4   (off-diagonal coupling, i ≠ j)"),
            B("hᵢ = −(Qᵢᵢ + Σⱼ Qᵢⱼ) / 4   (linear field from diagonal + row sum)"),
            E(),
            B("MaxCut: Wᵢⱼ = −wᵢⱼ for each edge, no bias — anti-aligned spins cut the edge"),
            B("MaxCut needs NO bias vector — simplest Ising embedding, ideal first demo"),
        ]
    ),
    (
        "Code Changes: hopfield_net.py → Ising Solver",
        [
            P("Three trivial edits unlock all combinatorial optimisation problems:"),
            E(),
            B("self.b = np.zeros(N)  →  add bias vector to __init__"),
            B("h_i = float(self.W[i] @ s) + self.b[i]  →  one-line change in run()"),
            B("energy() updated: −½ sᵀWs − bᵀs  →  correct for optimisation tracking"),
            E(),
            P("Plus three new methods:"),
            B("from_ising(J, h)  →  construct directly from Ising couplings"),
            B("from_qubo(Q)  →  QUBO → Ising → W, b automatically"),
            B("run_annealed(T_start, T_end, n_sweeps)  →  escape local minima via SA"),
        ]
    ),
    (
        "MaxCut — Simplest Ising Embedding",
        [
            P("Partition G=(V,E) into sets S and V\\S to maximise the weight of cut edges"),
            P("Assign sᵢ = +1 (in S) or −1 (not in S). Edge (i,j) is cut iff sᵢsⱼ = −1."),
            E(),
            B("Weight matrix: Wᵢⱼ = −wᵢⱼ for each edge — minimising E maximises the cut"),
            B("No bias vector needed — simplest problem to implement and verify"),
            E(),
            P("Benchmark datasets:"),
            B("Gset — 71 instances, 800–10,000 nodes (Stanford ORLIB, Burer et al. 2002)"),
            B("BiqMac Library — up to 3,000 nodes, known optimal solutions for verification"),
            B("MaxCutBench 2024 — standardised benchmark (Böther et al., arXiv:2406.11897)"),
        ]
    ),
    (
        "Travelling Salesman — Hopfield & Tank 1985",
        [
            P("Encode N cities × N positions as N² binary neurons V_{xi} (city x at position i in tour)"),
            E(),
            P("Energy has four penalty terms: row constraint + column constraint + total count + distance"),
            B("Penalty weights A,B,C,D must be tuned — too small → invalid tours, too large → random tours"),
            B("Only ~1–8% of runs yield valid tours for N=10 (Wilson & Pawley 1988)"),
            B("Weight matrix W is N²×N², bias b is length N² — far larger than memory-recall HNN"),
            E(),
            P("Benchmarks: TSPLIB (Reinelt 1991) — 100+ instances from N=51 to 1.8M cities"),
            P("Waterloo TSP data — known optimal solutions for correctness verification"),
        ]
    ),
    (
        "Other NP Problems as Ising / HNN",
        [
            P("Lucas (2014) gives explicit QUBO/Ising formulations for all of Karp's 21 NP-complete problems:"),
            E(),
            B("Number Partitioning — Wᵢⱼ = −2nᵢnⱼ, no bias. Second simplest after MaxCut"),
            B("Graph Coloring (k-color) — N·k spins, penalty per violated edge + per-vertex constraint"),
            B("Vertex Cover — N spins, penalty per uncovered edge (minimise |S| s.t. all edges covered)"),
            B("k-SAT — each clause adds a penalty; ground state = solution or nearest feasible state"),
            B("Hamiltonian Cycle, Clique, Set Cover — all map with O(N) to O(N²) spins"),
            E(),
            P("Reference: Lucas, A. (2014). Ising formulations of many NP problems. Frontiers in Physics 2:5"),
        ]
    ),
    (
        "Hardware Ising Machines — Landscape",
        [
            B("D-Wave (quantum annealing) — 5,627 superconducting qubits; sparse topology; 15 mK cryogenic"),
            B("Fujitsu Digital Annealer — ASIC, all-to-all connectivity, 8K–100K variables; parallel SA"),
            B("Toshiba SBM (FPGA) — Simulated Bifurcation Algorithm, 100K spins, deterministic, 0.5 ms on MaxCut"),
            B("Patel et al. Nature Electronics 2022 — FPGA p-bits via tanh LUT + LFSR; 10⁷× faster than D-Wave on MaxCut"),
            E(),
            P("Our LUT-HNN vs the field:"),
            B("Only approach using Boolean minimisation (Espresso SOP) of the update function — no arithmetic at inference"),
            B("Fully reconfigurable per problem instance — re-synthesise W for each new graph"),
            B("Deterministic updates (no LFSR noise) → add run_annealed() for escape"),
        ]
    ),
    (
        "Why Gradient-Based Training for HNNs?",
        [
            P("Classical closed-form rules saturate well below theoretical limits:"),
            B("Hebbian: 0.138N capacity — practical limit ~2 patterns for N=16"),
            B("Storkey: ~0.22N — better cross-talk correction, still far from maximum"),
            B("Pseudo-inverse: up to N exact fixed points, but offline and requires matrix inversion"),
            E(),
            P("Gradient-based methods can exceed these limits and support supervised objectives:"),
            B("Perceptron / 3-threshold SGD: ~0.64N capacity (Gardner bound) with local incremental updates"),
            B("Equilibrium Propagation: exact backprop gradient via two settle phases, no error wiring"),
            B("MPF (Sohl-Dickstein 2011): analytic gradient on single-bit-flip neighbours, ≥1 pattern/neuron"),
        ]
    ),
    (
        "Equilibrium Propagation (Scellier & Bengio 2017)",
        [
            P("Run network to free-phase fixed point s*, then nudge output units toward target → perturbed fixed point s^β"),
            E(),
            P("Weight update — purely local, Hebbian-style:"),
            B("ΔWᵢⱼ ∝ (sᵢ* sⱼ* − sᵢ^β sⱼ^β) / β"),
            B("Proven equivalent to BPTT via implicit function theorem on the fixed-point equation"),
            B("EqSpike (Laydevant 2021): spiking EP on neuromorphic silicon — 2–3 orders lower energy than GPU"),
            E(),
            P("FPGA relevance:"),
            B("Both phases reuse the same async settle circuit — no additional hardware for training"),
            B("Binary neurons: s^β − s* are sparse bit-flips → XOR-based weight increment logic"),
        ]
    ),
    (
        "SGD on the Fixed-Point Loss — Perceptron Rule",
        [
            P("Define a margin loss: pattern ξ^μ is stable iff every neuron i gets the right sign from the local field"),
            E(),
            P("L(W) = Σ_μ Σᵢ max(0,  −ξᵢ^μ · (Wᵢ · ξ^μ))"),
            E(),
            B("Gradient: ΔWᵢⱼ ∝ −Σ_{μ: violated} ξᵢ^μ ξⱼ^μ  — local, Hebbian, incremental"),
            B("Converges to Gardner bound: ~0.64N capacity (vs Hebbian 0.138N)"),
            B("Three-threshold variant (Alemi et al., PLoS Comp. Bio. 2015) adds learning/unlearning thresholds"),
            E(),
            P("FPGA: margin violation check is a single-bit comparison → simple on-chip FSM for online learning"),
        ]
    ),
    (
        "Training Method Comparison",
        [
            P("Method              |  Capacity   |  Supervised  |  FPGA Fit"),
            P("─────────────────────────────────────────────────────────────"),
            P("Hebbian             |  0.138N     |  No          |  Excellent"),
            P("Storkey             |  ~0.22N     |  No          |  Excellent"),
            P("Pseudo-inverse      |  N          |  No          |  Good (offline)"),
            P("Perceptron SGD      |  ~0.64N     |  No          |  Excellent (FSM)"),
            P("Equilibrium Prop.   |  Loss-dep   |  Yes         |  Good (EqSpike)"),
            P("MPF                 |  ≥1/neuron  |  No          |  Good (sparse)"),
            P("DEQ Adjoint         |  Loss-dep   |  Yes         |  Moderate"),
        ]
    ),
    (
        "Next Steps",
        [
            B("Add bias vector b to hopfield_net.py — unlocks all Ising / optimisation problems"),
            B("Implement from_ising(J, h) and from_qubo(Q) classmethods"),
            B("Demo: MaxCut on small graph — validate Ising mode vs known optimum"),
            B("Demo: Number partitioning — simplest bias-free problem after MaxCut"),
            E(),
            B("Implement perceptron / 3-threshold SGD rule — target ~0.64N capacity"),
            B("Explore Equilibrium Propagation for supervised pattern completion"),
            E(),
            B("Phase 2 hardware: integrate bias bᵢ into LUT truth table generation"),
            B("Espresso handles bᵢ transparently as a threshold shift — no pipeline changes"),
        ]
    ),
]

# ── open template and delete all slides ───────────────────────────────────────

prs = Presentation(INPUT)

# Delete all slides by removing from the slide list in presentation.xml
sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst):
    sldIdLst.remove(el)

# ── helper: add a slide by cloning a layout slide ─────────────────────────────

from pptx.util import Emu
from pptx.opc.packuri import PackURI
import pptx.oxml as oxml

def make_title_slide(prs, line1, line2):
    """Add a centered title slide (matches slide 1 layout from template)."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        ph.text_frame.clear()
    # Remove all shapes and rebuild
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:nvGrpSpPr') and sp.tag != qn('p:grpSpPr'):
            spTree.remove(sp)

    xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="2" name="Title"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph type="ctrTitle"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311708" y="1545450"/><a:ext cx="8520600" cy="2052600"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="b" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit fontScale="90000"/>
  </a:bodyPr>
  <a:lstStyle/>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="ctr">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{line1.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="ctr">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{line2.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(xml))
    return slide


def make_content_slide(prs, title_text, items):
    """Add a content slide with title + body (matches slide 3 layout)."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:nvGrpSpPr') and sp.tag != qn('p:grpSpPr'):
            spTree.remove(sp)

    # Title shape
    title_xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="2" name="Title"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph type="title"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311700" y="445025"/><a:ext cx="8520600" cy="572700"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="t" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit fontScale="90000"/>
  </a:bodyPr>
  <a:lstStyle/>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{title_text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(title_xml))

    # Body shape — build paragraphs
    paras_xml = ''
    for kind, text in items:
        if kind == 'p':
            paras_xml += plain_para(text)
        elif kind == 'b':
            paras_xml += bullet_para(text)
        else:
            paras_xml += empty_para()

    body_xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="3" name="Body"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph idx="1" type="body"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311700" y="1152475"/><a:ext cx="8520600" cy="3416400"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="t" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit/>
  </a:bodyPr>
  <a:lstStyle/>
  {paras_xml}
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(body_xml))
    return slide


# ── build slides ──────────────────────────────────────────────────────────────

make_title_slide(prs,
    "HNN as Ising Machines &",
    "Gradient-Based Training Methods")

for title, items in SLIDES[1:]:
    make_content_slide(prs, title, items)

prs.save(OUTPUT)
print(f"Saved → {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
