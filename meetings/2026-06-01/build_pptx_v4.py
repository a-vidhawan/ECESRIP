"""
Build SRIP_Presentation_v4.pptx
Topic: HNN as Cyclic Combinational Logic — Prior Work Deep Dive
Theme: cloned exactly from 0b712d9c-SRIP_Presena.pptx
"""

import copy, re, shutil
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

INPUT  = '/home/user/ECESRIP/meetings/2026-06-01/SRIP_Presentation_v3.pptx'
OUTPUT = '/home/user/ECESRIP/meetings/2026-06-01/SRIP_Presentation_v4.pptx'

A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── helpers ────────────────────────────────────────────────────────────────────

def _rpr(bold=False):
    return (f'<a:rPr xmlns:a="{A}" lang="en" sz="1730"'
            + (' b="1"' if bold else '')
            + f'><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:rPr>')

def _endParaRpr():
    return f'<a:endParaRPr xmlns:a="{A}" sz="1730"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:endParaRPr>'

def _esc(t):
    return t.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')

def plain_para(text, bold=False):
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buNone/></a:pPr>'
        f'<a:r>{_rpr(bold)}<a:t>{_esc(text)}</a:t></a:r>'
        f'{_endParaRpr()}</a:p>'
    )

def bullet_para(text):
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="-338455" lvl="0" marL="457200" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buClr><a:srgbClr val="000000"/></a:buClr>'
        f'<a:buSzPts val="1730"/><a:buChar char="&#x25CF;"/></a:pPr>'
        f'<a:r>{_rpr()}<a:t>{_esc(text)}</a:t></a:r>'
        f'{_endParaRpr()}</a:p>'
    )

def empty_para():
    return (
        f'<a:p xmlns:a="{A}">'
        f'<a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">'
        f'<a:lnSpc><a:spcPct val="95000"/></a:lnSpc>'
        f'<a:spcBef><a:spcPts val="0"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="0"/></a:spcAft>'
        f'<a:buNone/></a:pPr>'
        f'<a:r>{_rpr()}<a:t></a:t></a:r>'
        f'{_endParaRpr()}</a:p>'
    )

def P(label):  return ('p', label)
def B(label):  return ('b', label)
def E():       return ('e', '')

# ── slide content ─────────────────────────────────────────────────────────────

SLIDES = [
    (None, []),   # title — handled separately

    # 2. What Bill Asked — and the Finding
    (
        "What Bill Asked — and the Finding",
        [
            P("Task: find any prior work implementing HNN as a cyclic combinational circuit,"),
            P("with a Boolean threshold function precomputed at synthesis, not evaluated at runtime."),
            E(),
            P("Result after surveying 30+ papers across 40 years of HNN hardware:"),
            E(),
            B("No prior implementation precomputes the neuron update as a complete Boolean truth table"),
            B("No prior implementation uses direct feedback wires as recurrent connections"),
            B("Every digital implementation uses runtime multiply-accumulate + a digital comparator"),
            E(),
            P("The LUT-HNN is the first implementation that is structurally a cyclic combinational circuit."),
        ]
    ),

    # 3. How All Prior HNN Hardware Works
    (
        "How All Prior HNN Hardware Works",
        [
            P("The universal pattern — used in every analog, digital ASIC, FPGA, and memristor design:"),
            E(),
            B("Step 1: multiply each weight Wᵢⱼ by neighbour state sⱼ at runtime"),
            B("Step 2: accumulate the products to form the local field hᵢ = Σ Wᵢⱼ sⱼ"),
            B("Step 3: apply a digital or analog comparator: sᵢ ← sign(hᵢ)"),
            E(),
            P("The weight matrix W is always present in hardware:"),
            B("Analog: resistor ladders or CMOS synapse arrays (Graf 1987, Tank 1986)"),
            B("Digital FPGA: weights stored in BRAM, multiplied each clock cycle (Mansour 2011, Hikawa 2010)"),
            B("Memristor: crossbar current sum IS the weighted sum, comparator applied externally"),
        ]
    ),

    # 4. The LUT-HNN Difference
    (
        "The LUT-HNN Difference",
        [
            P("Instead of computing sᵢ ← sign(Σ Wᵢⱼ sⱼ) at runtime, precompute it entirely at synthesis:"),
            E(),
            B("For each neuron i, enumerate all 2^(N-1) possible input combinations of neighbour states"),
            B("Evaluate sign(Σ Wᵢⱼ sⱼ) for every combination — offline, once, at synthesis time"),
            B("Store the result as a Boolean truth table in the FPGA LUT fabric"),
            E(),
            P("At runtime:"),
            B("The N-1 neighbour states are wired directly as LUT address bits — no arithmetic"),
            B("The LUT output is sᵢ — the neuron's new state — with no comparator, no multiplier"),
            B("The weight matrix W does not exist at runtime. It was consumed by the LUT synthesis."),
            E(),
            P("This is a precomputed Boolean function, not a computed arithmetic expression."),
        ]
    ),

    # 5. Prior Work Survey: Who Built What
    (
        "Prior Work Survey: Who Built What",
        [
            P("Analog VLSI:"),
            B("Hopfield & Tank 1985 (Bell Labs, RC circuits + op-amps), Graf & Jackel 1987 (N=256 CMOS)"),
            B("Threshold = op-amp saturation. Weighted sum computed by parallel currents."),
            E(),
            P("Digital ASIC / FPGA:"),
            B("Jankowski et al. 1996, Hikawa 2003–10, Mansour 2011 — all: weight RAM + MAC + comparator"),
            B("p-Bit FPGA (Camsari 2017, Hassan DAC 2021) — tanh/sign LUT + LFSR noise + weight RAM + MAC"),
            E(),
            P("Emerging hardware:"),
            B("Memristor crossbars (Nature Comms 2015, 2021, 2024, 2025) — crossbar current = weighted sum"),
            B("Quantum HNN (IBM Q 2021) — majority vote over measurement shots"),
            B("Spin-memristor TLG (arXiv 2014) — device physics IS the threshold, gate-level only"),
        ]
    ),

    # 6. What "Cyclic Combinational" Means
    (
        "What \"Cyclic Combinational\" Means",
        [
            P("Standard digital logic teaching: feedback loops make a circuit sequential (stateful)."),
            P("Riedel & Bruck, DAC 2003 (Best Paper): this is not always true."),
            E(),
            B("A circuit with feedback is combinational if for every input, it has a unique stable output"),
            B("The output depends only on the current input, not on history — despite the feedback wires"),
            B("Such circuits can implement Boolean functions with fewer gates than any acyclic design"),
            E(),
            P("Well-behavedness (Riedel-Bruck): the circuit must converge to a unique fixed point"),
            P("for every input — no oscillation, no multi-stability for any given input."),
            E(),
            B("Cyclic circuits are not just theoretically allowed — they are synthesis-correct"),
            B("Tools (Espresso, ABC) can verify and optimise them; they map directly to FPGA LUT fabric"),
        ]
    ),

    # 7. The Bruck Thread
    (
        "The Bruck Thread: Convergence Theory to Cyclic Logic",
        [
            P("The same researcher — Jehoshua Bruck (Caltech) — authored both foundational papers:"),
            E(),
            B("Bruck 1990: 'On the Convergence Properties of the Hopfield Model,' Proc. IEEE 78(10)"),
            B("  Proves async HNN always converges — no input leads to sustained oscillation"),
            B("  This is precisely the well-behavedness condition for a cyclic combinational circuit"),
            E(),
            B("Riedel & Bruck 2003: 'The Synthesis of Cyclic Combinational Circuits,' DAC (Best Paper)"),
            B("  Formalises cyclic combinational logic — circuits with feedback that are still memoryless"),
            B("  Provides synthesis and verification framework; no connection to HNNs in the paper"),
            E(),
            P("13 years, same author, two frameworks that fit together exactly."),
            P("No paper has ever united them. The LUT-HNN is the first hardware realization of the union."),
        ]
    ),

    # 8. Formal Equivalence
    (
        "Formal Equivalence: Async HNN = Cyclic Combinational Circuit",
        [
            P("The mapping is exact, not approximate:"),
            E(),
            B("Each neuron i implements a Boolean threshold function fᵢ(s₁,...,sₙ) = sign(Σ Wᵢⱼ sⱼ)"),
            B("The recurrent connections Wᵢⱼ become feedback wires: output of neuron j feeds input of neuron i"),
            B("The network has feedback loops — it is cyclic by construction"),
            E(),
            P("Why it is combinational (not sequential):"),
            B("Bruck 1990 proves async update always converges to a fixed point — every input is well-behaved"),
            B("At a fixed point, fᵢ(s*) = sᵢ* for all i — outputs are consistent with inputs"),
            B("The stable output depends only on initial state, not on update history"),
            E(),
            P("This satisfies the Riedel-Bruck definition exactly. The HNN IS a cyclic combinational circuit."),
        ]
    ),

    # 9. Two Variants
    (
        "Two Implementation Variants",
        [
            P("Phase 1 — Clocked LUT (implemented):"),
            B("Flip-flops hold state. Each clock edge: all LUT inputs sampled, outputs registered."),
            B("Update is synchronous but logically identical to async — one neuron updated per cycle"),
            B("Standard FPGA design flow; timing closure straightforward"),
            E(),
            P("Phase 2 — Combinational Feedback (proposed, novel):"),
            B("No flip-flops. LUT outputs wire directly back to LUT inputs."),
            B("Settlement occurs through propagation delay — purely combinational relaxation"),
            B("No clock signal at all. The circuit IS the Riedel-Bruck cyclic combinational realization."),
            B("Convergence guaranteed by Bruck 1990 — propagation delay acts as the async update schedule"),
            E(),
            P("Phase 2 has never been built. It is the cleanest instantiation of the theory."),
        ]
    ),

    # 10. Closest Prior Work and Why It's Different
    (
        "Closest Prior Work: p-Bit FPGA — Still Different",
        [
            P("The closest prior work: Camsari et al. (Physical Review X 2017, arXiv 2017) and"),
            P("Hassan et al. (DAC 2021, arXiv:2101.00147) — FPGA implementations of p-Bits."),
            E(),
            P("What they do that looks similar:"),
            B("Use a LUT on FPGA for the activation function (tanh approximation)"),
            B("Async sequential update order — one neuron at a time"),
            E(),
            P("Why they are fundamentally different:"),
            B("The LUT stores a piecewise-linear approximation of tanh(x) — a function of a scalar input"),
            B("The weighted sum Σ Wᵢⱼ sⱼ is still computed at runtime; W is still in BRAM"),
            B("The LUT is a function approximator, not a truth table over all 2^(N-1) inputs"),
            B("An LFSR generates stochastic noise; the update is probabilistic, not deterministic"),
            E(),
            P("The LUT-HNN uses no W at runtime, no MAC, no noise — just address-in, bit-out."),
        ]
    ),

    # 11. Novelty Claim
    (
        "Novelty Claim",
        [
            P("The LUT-HNN is the first hardware implementation that:"),
            E(),
            B("Precomputes the complete Boolean truth table of each neuron's threshold function"),
            B("Eliminates all runtime arithmetic — no multipliers, no adders, no comparators"),
            B("Uses FPGA feedback routing as the literal recurrent weight connections"),
            B("Structurally instantiates the Riedel-Bruck cyclic combinational circuit framework"),
            E(),
            P("The theoretical foundations have existed since 1990 (Bruck) and 2003 (Riedel-Bruck)."),
            P("No prior implementation — analog, digital ASIC, FPGA, memristor, or quantum — has"),
            P("used this architecture. The architectural gap is confirmed across 40 years of literature."),
            E(),
            B("Next: verify Phase 2 (combinational feedback) in simulation; scale to N=20 benchmark"),
        ]
    ),
]

# ── open template ──────────────────────────────────────────────────────────────

prs = Presentation(INPUT)

sldIdLst = prs.slides._sldIdLst
for el in list(sldIdLst):
    sldIdLst.remove(el)

# ── slide builders ─────────────────────────────────────────────────────────────

from pptx.util import Emu
from pptx.opc.packuri import PackURI
import pptx.oxml as oxml


def make_title_slide(prs, line1, line2):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:nvGrpSpPr') and sp.tag != qn('p:grpSpPr'):
            spTree.remove(sp)

    xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="2" name="Title"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph type="ctrTitle"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311708" y="1545450"/><a:ext cx="8520600" cy="2052600"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="b" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit fontScale="90000"/>
  </a:bodyPr>
  <a:lstStyle/>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="ctr">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{_esc(line1)}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="ctr">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{_esc(line2)}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(xml))
    return slide


def make_content_slide(prs, title_text, items):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:nvGrpSpPr') and sp.tag != qn('p:grpSpPr'):
            spTree.remove(sp)

    title_xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="2" name="Title"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph type="title"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311700" y="445025"/><a:ext cx="8520600" cy="572700"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="t" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit fontScale="90000"/>
  </a:bodyPr>
  <a:lstStyle/>
  <a:p>
    <a:pPr indent="0" lvl="0" marL="0" rtl="0" algn="l">
      <a:spcBef><a:spcPts val="0"/></a:spcBef>
      <a:spcAft><a:spcPts val="0"/></a:spcAft>
      <a:buNone/>
    </a:pPr>
    <a:r><a:rPr lang="en"/><a:t>{_esc(title_text)}</a:t></a:r>
    <a:endParaRPr/>
  </a:p>
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(title_xml))

    paras_xml = ''
    for kind, text in items:
        if kind == 'p':
            paras_xml += plain_para(text)
        elif kind == 'b':
            paras_xml += bullet_para(text)
        else:
            paras_xml += empty_para()

    body_xml = f'''<p:sp xmlns:a="{A}" xmlns:p="{PML}">
<p:nvSpPr>
  <p:cNvPr id="3" name="Body"/>
  <p:cNvSpPr txBox="1"/>
  <p:nvPr><p:ph idx="1" type="body"/></p:nvPr>
</p:nvSpPr>
<p:spPr>
  <a:xfrm><a:off x="311700" y="1152475"/><a:ext cx="8520600" cy="3416400"/></a:xfrm>
  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
</p:spPr>
<p:txBody>
  <a:bodyPr anchorCtr="0" anchor="t" bIns="91425" lIns="91425"
            spcFirstLastPara="1" rIns="91425" wrap="square" tIns="91425">
    <a:normAutofit/>
  </a:bodyPr>
  <a:lstStyle/>
  {paras_xml}
</p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(body_xml))
    return slide


# ── build ──────────────────────────────────────────────────────────────────────

make_title_slide(prs,
    "HNN as Cyclic Combinational Logic",
    "Prior Work Deep Dive")

for title, items in SLIDES[1:]:
    make_content_slide(prs, title, items)

prs.save(OUTPUT)
print(f"Saved -> {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
